#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from pypdf import PdfReader
from pptx import Presentation


def clean_text(text: str) -> str:
    text = text.replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    chunks = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")
    return clean_text("\n\n".join(chunks))


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slide_chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        text_runs: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
        slide_text = clean_text("\n".join(text_runs))
        if slide_text:
            slide_chunks.append(f"[Slide {i}]\n{slide_text}")
    return clean_text("\n\n".join(slide_chunks))


def extract_docx(path: Path) -> str:
    import zipfile
    from xml.etree import ElementTree as ET

    with zipfile.ZipFile(path) as zf:
        data = zf.read("word/document.xml")
    root = ET.fromstring(data)
    texts = []
    for node in root.iter():
        if node.tag.endswith("}t") and node.text:
            texts.append(node.text)
    return clean_text("\n".join(texts))


def extract_tex(path: Path) -> str:
    return clean_text(path.read_text(encoding="utf-8", errors="replace"))


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".tex":
        return extract_tex(path)
    raise ValueError(f"Unsupported extension: {suffix}")


def load_inventory(path: Path) -> list[dict[str, Any]]:
    records = []
    for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
        if not line.strip():
            continue
        obj = json.loads(line)
        if obj.get("type") == "meta":
            continue
        records.append(obj)
    return records


def write_json(path: Path, obj: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(obj, indent=2, ensure_ascii=True) + "\n", encoding="utf-8")


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract text from inventoried papers and decks.")
    ap.add_argument("--inventory", required=True)
    ap.add_argument("--out-dir", required=True)
    ap.add_argument("--limit", type=int, default=0, help="Optional max documents to extract")
    ap.add_argument("--kind", default="", help="Optional kind filter")
    args = ap.parse_args()

    records = load_inventory(Path(args.inventory))
    out_dir = Path(args.out_dir)
    selected = []
    for rec in records:
        if args.kind and rec.get("kind") != args.kind:
            continue
        selected.append(rec)
    if args.limit > 0:
        selected = selected[: args.limit]

    manifest_rows = []
    for rec in selected:
        path = Path(rec["path"])
        try:
            text = extract_text(path)
            status = "ok"
            error = ""
        except Exception as exc:
            text = ""
            status = "error"
            error = str(exc)

        snippet = text[:2000]
        out_path = out_dir / f"{rec['doc_id']}.json"
        payload = {
            "doc_id": rec["doc_id"],
            "path": rec["path"],
            "title_guess": rec["title_guess"],
            "kind": rec["kind"],
            "collection_id": rec["collection_id"],
            "status": status,
            "error": error,
            "char_count": len(text),
            "snippet": snippet,
            "text": text,
        }
        write_json(out_path, payload)
        manifest_rows.append(
            {
                "doc_id": rec["doc_id"],
                "status": status,
                "error": error,
                "char_count": len(text),
                "path": rec["path"],
                "output": str(out_path),
            }
        )

    manifest_path = out_dir / "manifest.json"
    write_json(manifest_path, {"count": len(manifest_rows), "rows": manifest_rows})
    print(f"extracted={len(manifest_rows)}")
    print(f"manifest={manifest_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
