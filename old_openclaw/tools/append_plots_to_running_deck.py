#!/usr/bin/env python3
"""Append one slide with up to 2 images to an existing running deck.

Usage:
  .venv_pptx/bin/python tools/append_plots_to_running_deck.py \
    --pptx out/ppt/running/.../deck.pptx \
    --title "..." \
    --img1 out/plots/...png --img2 out/plots/...png

Non-destructive-ish: writes in-place (makes a .bak copy).
"""

from __future__ import annotations

import argparse
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _add_title(slide, title: str):
    # top title box
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(12.3)
    height = Inches(0.6)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--title", required=True)
    ap.add_argument("--img1", required=True)
    ap.add_argument("--img2")
    args = ap.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        raise SystemExit(f"PPTX not found: {pptx_path}")

    # backup
    bak = pptx_path.with_suffix(pptx_path.suffix + ".bak")
    if not bak.exists():
        shutil.copy2(pptx_path, bak)

    prs = Presentation(str(pptx_path))

    # choose blank layout when possible
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    _add_title(slide, args.title)

    # image placement (16:9-ish decks; python-pptx uses inches)
    # two-up layout
    top = Inches(1.0)
    height = Inches(5.8)
    if args.img2:
        slide.shapes.add_picture(str(Path(args.img1)), Inches(0.6), top, width=Inches(6.2), height=height)
        slide.shapes.add_picture(str(Path(args.img2)), Inches(6.9), top, width=Inches(6.2), height=height)
    else:
        slide.shapes.add_picture(str(Path(args.img1)), Inches(0.8), top, width=Inches(12.0), height=height)

    prs.save(str(pptx_path))


if __name__ == "__main__":
    main()
