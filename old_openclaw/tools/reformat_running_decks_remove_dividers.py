"""Remove placeholder section-divider slides from OpenScience running decks.

Rationale: early versions of the running deck template included several divider slides
with only a title (e.g., "Project framing", "Current best-known results"). These read
as unused/blank in PI-facing sharing.

This script removes those divider slides (by matching exact title text) and rewrites
PPTX in-place.

Safety:
- Designed for derived artifacts under out/ppt/** only.
- Make backups before running.

Usage:
  . .venv_pptx/bin/activate
  python tools/reformat_running_decks_remove_dividers.py out/ppt/running/**/**.pptx

"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


DIVIDER_TITLES = {
    "Project framing (device, modality)",
    "Current best-known results",
    "New analyses (rolling)",
    "Open questions / next experiments",
    "Provenance / notes",
}


def slide_all_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            t = (para.text or "").strip()
            if t:
                texts.append(t)
    return "\n".join(texts).strip()


def delete_slide(prs: Presentation, slide_idx0: int) -> None:
    """Delete slide at 0-based index from python-pptx Presentation.

    Implementation uses private API patterns commonly used with python-pptx.
    """
    slide_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
    slide_id_elems = list(slide_id_list)
    slide_id = slide_id_elems[slide_idx0]

    rId = slide_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    # drop from slide id list
    slide_id_list.remove(slide_id)
    # drop relationship
    prs.part.drop_rel(rId)


def process(path: Path) -> dict:
    prs = Presentation(str(path))
    to_delete = []
    for i, slide in enumerate(prs.slides):
        txt = slide_all_text(slide)
        # If the entire slide's text is exactly one of the divider titles, it's a placeholder.
        if txt in DIVIDER_TITLES:
            to_delete.append(i)

    # delete in reverse order so indices stay valid
    for i in reversed(to_delete):
        delete_slide(prs, i)

    prs.save(str(path))
    return {"path": str(path), "slides_deleted": len(to_delete)}


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print("Provide one or more pptx paths", file=sys.stderr)
        return 2

    paths = []
    for a in argv[1:]:
        p = Path(a)
        if any(ch in a for ch in "*?["):
            # glob
            paths.extend([pp for pp in Path().glob(a) if pp.suffix.lower() == ".pptx"])
        else:
            paths.append(p)

    ok = True
    for p in paths:
        if not p.exists() or p.suffix.lower() != ".pptx":
            print(f"[skip] {p}")
            ok = False
            continue
        res = process(p)
        print(f"[done] {res['path']}: deleted {res['slides_deleted']} divider slides")

    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
